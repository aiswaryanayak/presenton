# servers/fastapi/services/pptx_presentation_creator.py
import os
import uuid
import asyncio
import json
from typing import List, Optional, Any

from lxml import etree
from lxml.etree import fromstring, tostring
from PIL import Image

from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.slide import Slide
from pptx.text.text import _Paragraph, TextFrame, Font, _Run
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt
from pptx.dml.color import RGBColor

# Local helpers (make sure these paths match your repo)
from services.html_to_text_runs_service import parse_html_text_to_text_runs as parse_inline_html_to_runs
from utils.download_helpers import download_files
from utils.image_utils import (
    clip_image,
    create_circle_image,
    fit_image,
    invert_image,
    round_image_corners,
    set_image_opacity,
)

# Pptx models from the file you pasted
from models.pptx_models import (
    PptxAutoShapeBoxModel,
    PptxBoxShapeEnum,
    PptxConnectorModel,
    PptxFillModel,
    PptxFontModel,
    PptxParagraphModel,
    PptxPictureBoxModel,
    PptxPositionModel,
    PptxPresentationModel,
    PptxShadowModel,
    PptxSlideModel,
    PptxSpacingModel,
    PptxStrokeModel,
    PptxTextBoxModel,
    PptxTextRunModel,
    PptxPictureModel,
)

# index of a blank layout in the default template; your environment may vary
BLANK_SLIDE_LAYOUT = 6


class PptxPresentationCreator:
    def __init__(self, ppt_model: PptxPresentationModel, temp_dir: str):
        self._temp_dir = temp_dir
        self._ppt_model = ppt_model
        self._slide_models: List[PptxSlideModel] = getattr(ppt_model, "slides", []) or []

        self._ppt = Presentation()
        # keep the size you used previously
        self._ppt.slide_width = Pt(1280)
        self._ppt.slide_height = Pt(720)

    def get_sub_element(self, parent, tagname, **kwargs):
        el = OxmlElement(tagname)
        for k, v in kwargs.items():
            el.set(k, str(v))
        parent.append(el)
        return el

    async def fetch_network_assets(self):
        """
        Download http(s) images referenced in slides into the temp dir.
        Update picture.path to local file path when downloaded.
        """
        image_urls: List[str] = []
        models_with_network_asset: List[PptxPictureBoxModel] = []

        # gather global shapes if present
        global_shapes = getattr(self._ppt_model, "shapes", []) or []
        for s in global_shapes:
            if isinstance(s, PptxPictureBoxModel):
                pic_path = getattr(s.picture, "path", "")
                if isinstance(pic_path, str) and pic_path.startswith("http"):
                    image_urls.append(pic_path)
                    models_with_network_asset.append(s)

        # gather per-slide shapes
        for slide in self._slide_models:
            for shape in getattr(slide, "shapes", []) or []:
                if isinstance(shape, PptxPictureBoxModel):
                    pic_path = getattr(shape.picture, "path", "")
                    if isinstance(pic_path, str) and pic_path.startswith("http"):
                        image_urls.append(pic_path)
                        models_with_network_asset.append(shape)

        if not image_urls:
            return

        try:
            downloaded_paths = await download_files(image_urls, self._temp_dir)
        except Exception as e:
            print("❌ download_files failed:", e)
            downloaded_paths = [None] * len(image_urls)

        for model, local_path in zip(models_with_network_asset, downloaded_paths):
            if local_path:
                model.picture.path = local_path
                model.picture.is_network = False
            else:
                model.picture.path = "/static/images/placeholder.jpg"
                model.picture.is_network = False

    async def create_ppt(self):
        await self.fetch_network_assets()

        for slide_model in self._slide_models:
            # Build shapes list - do not mutate global shapes list
            global_shapes = getattr(self._ppt_model, "shapes", []) or []
            slide_shapes = list(getattr(slide_model, "shapes", []) or [])
            if isinstance(global_shapes, list) and global_shapes:
                # Prepend global shapes so they render beneath slide-specific shapes
                slide_shapes = list(global_shapes) + slide_shapes

            # We'll iterate the resolved shapes list (don't try to reassign into pydantic models)
            self.add_and_populate_slide(slide_model, slide_shapes)

    def set_presentation_theme(self):
        """
        Optionally set theme colors if the model provides them (model may include a 'theme' dict).
        """
        try:
            theme = getattr(self._ppt_model, "theme", {}) or {}
            colors = theme.get("colors", {}) or {}
            if not colors:
                return
            slide_master = self._ppt.slide_master
            slide_master_part = slide_master.part
            theme_part = slide_master_part.part_related_by(RT.THEME)
            xml = fromstring(theme_part.blob)
            nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
            for name, hex_value in colors.items():
                nodes = xml.xpath(f"a:themeElements/a:clrScheme/a:{name}/a:srgbClr", namespaces=nsmap)
                if nodes:
                    node = nodes[0]
                    node.set("val", hex_value.replace("#", "").upper())
            theme_part._blob = tostring(xml)
        except Exception as e:
            print("⚠️ set_presentation_theme failed:", e)

    def add_and_populate_slide(self, slide_model: PptxSlideModel, shapes_override: Optional[List[Any]] = None):
        """
        Create a slide and populate it with shapes.
        shapes_override: list of shape models to use (already resolved with globals)
        """
        try:
            slide = self._ppt.slides.add_slide(self._ppt.slide_layouts[BLANK_SLIDE_LAYOUT])
        except Exception:
            slide = self._ppt.slides.add_slide(self._ppt.slide_layouts[0])

        if getattr(slide_model, "background", None):
            try:
                self.apply_fill_to_shape(slide.background, slide_model.background)
            except Exception as e:
                print("⚠️ apply_fill failed on background:", e)

        if getattr(slide_model, "note", None):
            try:
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    slide.notes_slide.notes_text_frame.text = slide_model.note
            except Exception:
                pass

        shapes_to_add = shapes_override if shapes_override is not None else list(getattr(slide_model, "shapes", []) or [])

        for shape_model in shapes_to_add:
            try:
                if isinstance(shape_model, PptxPictureBoxModel):
                    self.add_picture(slide, shape_model)
                elif isinstance(shape_model, PptxAutoShapeBoxModel):
                    self.add_autoshape(slide, shape_model)
                elif isinstance(shape_model, PptxTextBoxModel):
                    self.add_textbox(slide, shape_model)
                elif isinstance(shape_model, PptxConnectorModel):
                    self.add_connector(slide, shape_model)
                else:
                    # fallback: add string representation as a textbox
                    try:
                        fallback_txt = str(shape_model)
                        pos = PptxPositionModel(left=50, top=50, width=500, height=100)
                        fallback_para = PptxParagraphModel(text=fallback_txt, font=PptxFontModel())
                        fallback_box = PptxTextBoxModel(position=pos, paragraphs=[fallback_para])
                        self.add_textbox(slide, fallback_box)
                    except Exception:
                        print("⚠️ Unknown shape type and fallback failed:", type(shape_model))
            except Exception as e:
                print("⚠️ adding shape failed:", e)

    def add_connector(self, slide: Slide, connector_model: PptxConnectorModel):
        try:
            if getattr(connector_model, "thickness", 0) == 0:
                return
            connector_shape = slide.shapes.add_connector(
                connector_model.type, *connector_model.position.to_pt_xyxy()
            )
            connector_shape.line.width = Pt(connector_model.thickness)
            connector_shape.line.color.rgb = RGBColor.from_string(connector_model.color)
            self.set_fill_opacity(connector_shape, connector_model.opacity)
        except Exception as e:
            print("⚠️ add_connector failed:", e)

    def add_picture(self, slide: Slide, picture_model: PptxPictureBoxModel):
        img_path = getattr(picture_model.picture, "path", None)
        if not img_path:
            return

        # process transformations if requested
        try:
            needs_processing = bool(
                getattr(picture_model, "clip", False)
                or getattr(picture_model, "border_radius", None)
                or getattr(picture_model, "invert", False)
                or getattr(picture_model, "opacity", None)
                or getattr(picture_model, "object_fit", None)
                or getattr(picture_model, "shape", None)
            )
            if needs_processing:
                image = Image.open(img_path).convert("RGBA")

                # border_radius might be a list or int in your model; take first element if list
                br = getattr(picture_model, "border_radius", None)
                br_val = None
                if isinstance(br, list) and br:
                    try:
                        br_val = int(br[0])
                    except Exception:
                        br_val = None
                elif isinstance(br, int):
                    br_val = br

                if br_val:
                    image = round_image_corners(image, br_val)

                if getattr(picture_model, "object_fit", None):
                    # object_fit is a PptxObjectFitModel - we only honor fit type
                    fit = getattr(picture_model.object_fit, "fit", None)
                    if fit:
                        image = fit_image(
                            image,
                            picture_model.position.width,
                            picture_model.position.height,
                            fit.value if hasattr(fit, "value") else str(fit),
                        )
                elif getattr(picture_model, "clip", False):
                    image = clip_image(
                        image,
                        picture_model.position.width,
                        picture_model.position.height,
                    )

                if getattr(picture_model, "shape", None) == PptxBoxShapeEnum.CIRCLE:
                    image = create_circle_image(image)

                if getattr(picture_model, "invert", False):
                    image = invert_image(image)

                if getattr(picture_model, "opacity", None) is not None:
                    image = set_image_opacity(image, picture_model.opacity)

                out_path = os.path.join(self._temp_dir, f"{uuid.uuid4()}.png")
                image.save(out_path)
                img_path = out_path
        except Exception as e:
            print("⚠️ image processing failed:", e)
            if not os.path.exists(img_path):
                img_path = "/static/images/placeholder.jpg"

        try:
            margined = self.get_margined_position(picture_model.position, getattr(picture_model, "margin", None))
            slide.shapes.add_picture(img_path, *margined.to_pt_list())
        except Exception as e:
            print("⚠️ add_picture to slide failed:", e)

    def add_autoshape(self, slide: Slide, autoshape_box_model: PptxAutoShapeBoxModel):
        try:
            position = autoshape_box_model.position
            if getattr(autoshape_box_model, "margin", None):
                position = self.get_margined_position(position, autoshape_box_model.margin)

            autoshape = slide.shapes.add_shape(
                autoshape_box_model.type, *position.to_pt_list()
            )

            textbox = autoshape.text_frame
            textbox.word_wrap = autoshape_box_model.text_wrap

            self.apply_fill_to_shape(autoshape, autoshape_box_model.fill)
            self.apply_margin_to_text_box(textbox, autoshape_box_model.margin)
            self.apply_stroke_to_shape(autoshape, autoshape_box_model.stroke)
            self.apply_shadow_to_shape(autoshape, autoshape_box_model.shadow)
            self.apply_border_radius_to_shape(autoshape, autoshape_box_model.border_radius)

            if getattr(autoshape_box_model, "paragraphs", None):
                self.add_paragraphs(textbox, autoshape_box_model.paragraphs)
        except Exception as e:
            print("⚠️ add_autoshape failed:", e)

    def add_textbox(self, slide: Slide, textbox_model: PptxTextBoxModel):
        try:
            position = textbox_model.position
            if getattr(textbox_model, "margin", None):
                position = self.get_margined_position(position, textbox_model.margin)

            textbox_shape = slide.shapes.add_textbox(*position.to_pt_list())
            textbox_shape.width += Pt(2)
            textbox = textbox_shape.text_frame
            textbox.word_wrap = textbox_model.text_wrap

            self.apply_fill_to_shape(textbox_shape, textbox_model.fill)
            self.apply_margin_to_text_box(textbox, textbox_model.margin)
            self.add_paragraphs(textbox, textbox_model.paragraphs)
        except Exception as e:
            print("⚠️ add_textbox failed:", e)

    def add_paragraphs(self, textbox: TextFrame, paragraph_models: List[PptxParagraphModel]):
        for idx, para_model in enumerate(paragraph_models):
            paragraph = textbox.add_paragraph() if idx > 0 else textbox.paragraphs[0]
            self.populate_paragraph(paragraph, para_model)

    def populate_paragraph(self, paragraph: _Paragraph, paragraph_model: PptxParagraphModel):
        try:
            if paragraph_model.spacing:
                self.apply_spacing_to_paragraph(paragraph, paragraph_model.spacing)

            if paragraph_model.line_height:
                paragraph.line_spacing = paragraph_model.line_height

            if paragraph_model.alignment is not None:
                paragraph.alignment = paragraph_model.alignment

            if paragraph_model.font:
                self.apply_font_to_paragraph(paragraph, paragraph_model.font)

            text_runs = []
            if paragraph_model.text:
                text_runs = self.parse_html_text_to_text_runs(paragraph_model.font, paragraph_model.text)
            elif paragraph_model.text_runs:
                text_runs = paragraph_model.text_runs

            for run_model in text_runs:
                text_run = paragraph.add_run()
                self.populate_text_run(text_run, run_model)
        except Exception as e:
            print("⚠️ populate_paragraph failed:", e)

    def parse_html_text_to_text_runs(self, font: Optional[PptxFontModel], text: str):
        return parse_inline_html_to_runs(text, font)

    def populate_text_run(self, text_run: _Run, text_run_model: PptxTextRunModel):
        try:
            text_run.text = text_run_model.text
            if text_run_model.font:
                self.apply_font(text_run.font, text_run_model.font)
        except Exception as e:
            print("⚠️ populate_text_run failed:", e)

    def apply_border_radius_to_shape(self, shape: Shape, border_radius: Optional[int]):
        if not border_radius:
            return
        try:
            normalized = Pt(border_radius) / max(min(shape.width, shape.height), Pt(1))
            try:
                # some shapes don't support adjustments
                shape.adjustments[0] = normalized
            except Exception:
                pass
        except Exception as e:
            print("⚠️ apply_border_radius_to_shape failed:", e)

    def apply_fill_to_shape(self, shape: Shape, fill: Optional[PptxFillModel] = None):
        try:
            if not fill:
                shape.fill.background()
                return
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(fill.color)
            self.set_fill_opacity(shape.fill, fill.opacity)
        except Exception as e:
            print("⚠️ apply_fill_to_shape failed:", e)

    def apply_stroke_to_shape(self, shape: Shape, stroke: Optional[PptxStrokeModel] = None):
        try:
            if not stroke or getattr(stroke, "thickness", 0) == 0:
                shape.line.fill.background()
                return
            shape.line.fill.solid()
            shape.line.fill.fore_color.rgb = RGBColor.from_string(stroke.color)
            shape.line.width = Pt(stroke.thickness)
            self.set_fill_opacity(shape.line.fill, stroke.opacity)
        except Exception as e:
            print("⚠️ apply_stroke_to_shape failed:", e)

    def apply_shadow_to_shape(self, shape: Shape, shadow: Optional[PptxShadowModel] = None):
        try:
            sp_element = shape._element
            sp_pr = sp_element.xpath("p:spPr")[0]
            nsmap = sp_pr.nsmap or {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
            effect_list = sp_pr.find("a:effectLst", namespaces=nsmap)
            if effect_list is None:
                effect_list = etree.SubElement(sp_pr, f"{{{nsmap['a']}}}effectLst")

            old = effect_list.find("a:outerShdw", namespaces=nsmap)
            if old is not None:
                effect_list.remove(old)

            if shadow is None:
                outer = etree.SubElement(effect_list, f"{{{nsmap['a']}}}outerShdw", {"blurRad": "0", "dist": "0", "dir": "0"})
                color_el = etree.SubElement(outer, f"{{{nsmap['a']}}}srgbClr", {"val": "000000"})
                etree.SubElement(color_el, f"{{{nsmap['a']}}}alpha", {"val": "0"})
                return

            blur = str(int(Pt(shadow.radius))) if getattr(shadow, "radius", None) else "40000"
            dist = str(int(Pt(shadow.offset))) if getattr(shadow, "offset", None) else "40000"
            dir_val = str(int(getattr(shadow, "angle", 0) * 1000))
            outer = etree.SubElement(effect_list, f"{{{nsmap['a']}}}outerShdw", {"blurRad": blur, "dir": dir_val, "dist": dist, "rotWithShape": "0"})
            color_el = etree.SubElement(outer, f"{{{nsmap['a']}}}srgbClr", {"val": getattr(shadow, "color", "000000")})
            alpha_val = str(int(getattr(shadow, "opacity", 0.0) * 100000)) if getattr(shadow, "opacity", None) is not None else "100000"
            etree.SubElement(color_el, f"{{{nsmap['a']}}}alpha", {"val": alpha_val})
        except Exception as e:
            print("⚠️ apply_shadow_to_shape failed:", e)

    def set_fill_opacity(self, fill, opacity):
        if opacity is None or opacity >= 1.0:
            return
        try:
            alpha = int(opacity * 100000)
            ts = getattr(fill, "_xPr", None)
            if not ts:
                return
            # best-effort: insert alpha element
            try:
                sF = ts.solidFill if hasattr(ts, "solidFill") else None
                if sF is None:
                    sF = ts.get_or_change_to_srgbClr()
                self.get_sub_element(sF, "a:alpha", val=str(alpha))
            except Exception:
                pass
        except Exception as e:
            print("⚠️ set_fill_opacity failed:", e)

    def get_margined_position(self, position: PptxPositionModel, margin: Optional[PptxSpacingModel]) -> PptxPositionModel:
        if not margin:
            return position
        left = position.left + margin.left
        top = position.top + margin.top
        width = max(position.width - margin.left - margin.right, 0)
        height = max(position.height - margin.top - margin.bottom, 0)
        return PptxPositionModel(left=left, top=top, width=width, height=height)

    def apply_margin_to_text_box(self, text_frame: TextFrame, margin: Optional[PptxSpacingModel]):
        try:
            text_frame.margin_left = Pt(margin.left if margin else 0)
            text_frame.margin_right = Pt(margin.right if margin else 0)
            text_frame.margin_top = Pt(margin.top if margin else 0)
            text_frame.margin_bottom = Pt(margin.bottom if margin else 0)
        except Exception as e:
            print("⚠️ apply_margin_to_text_box failed:", e)

    def apply_spacing_to_paragraph(self, paragraph: _Paragraph, spacing: PptxSpacingModel):
        try:
            paragraph.space_before = Pt(spacing.top)
            paragraph.space_after = Pt(spacing.bottom)
        except Exception as e:
            print("⚠️ apply_spacing_to_paragraph failed:", e)

    def apply_font_to_paragraph(self, paragraph: _Paragraph, font: PptxFontModel):
        try:
            self.apply_font(paragraph.font, font)
        except Exception as e:
            print("⚠️ apply_font_to_paragraph failed:", e)

    def apply_font(self, font: Font, font_model: PptxFontModel):
        try:
            if getattr(font_model, "name", None):
                font.name = font_model.name
            if getattr(font_model, "color", None):
                font.color.rgb = RGBColor.from_string(font_model.color)
            font.italic = bool(getattr(font_model, "italic", False))
            if getattr(font_model, "size", None):
                font.size = Pt(font_model.size)
            font.bold = bool(getattr(font_model, "font_weight", 400) >= 600)
            if getattr(font_model, "underline", None) is not None:
                font.underline = bool(font_model.underline)
            if getattr(font_model, "strike", None) is not None:
                self.apply_strike_to_font(font, font_model.strike)
        except Exception as e:
            print("⚠️ apply_font failed:", e)

    def apply_strike_to_font(self, font: Font, strike: Optional[bool]):
        try:
            rPr = font._element
            if strike is True:
                rPr.set("strike", "sngStrike")
            elif strike is False:
                rPr.set("strike", "noStrike")
        except Exception as e:
            print("⚠️ apply_strike_to_font failed:", e)

    def save(self, path: str):
        os.makedirs(os.path.dirname(path), exist_ok=True)
        try:
            self._ppt.save(path)
        except Exception as e:
            print("❌ save failed:", e)
            raise
