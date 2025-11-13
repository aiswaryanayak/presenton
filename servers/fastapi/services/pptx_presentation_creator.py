# servers/fastapi/services/pptx_presentation_creator.py
import os
import uuid
import json
import asyncio
from typing import List, Optional

from lxml import etree
from lxml.etree import fromstring, tostring
from PIL import Image

from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.slide import Slide
from pptx.text.text import _Paragraph, TextFrame, Font, _Run
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt
from pptx.dml.color import RGBColor

# Local imports (models & helpers)
from services.html_to_text_runs_service import (
    parse_html_text_to_text_runs as parse_inline_html_to_runs,
)
from utils.download_helpers import download_files
from utils.image_utils import (
    clip_image,
    create_circle_image,
    fit_image,
    invert_image,
    round_image_corners,
    set_image_opacity,
)

# Import Pptx models. Keep relative import to your actual models module.
from models.pptx_models import (
    PptxAutoShapeBoxModel,
    PptxBoxShapeEnum,
    PptxConnectorModel,
    PptxFillModel,
    PptxFontModel,
    PptxParagraphModel,
    PptxPictureBoxModel,
    PptxPositionModel,
    PptxPresentationModel,
    PptxShadowModel,
    PptxSlideModel,
    PptxSpacingModel,
    PptxStrokeModel,
    PptxTextBoxModel,
    PptxTextRunModel,
)

BLANK_SLIDE_LAYOUT = 6  # keep as before; ensure PowerPoint has this layout index


class PptxPresentationCreator:
    def __init__(self, ppt_model: PptxPresentationModel, temp_dir: str):
        self._temp_dir = temp_dir
        self._ppt_model = ppt_model
        self._slide_models = getattr(ppt_model, "slides", []) or []

        self._ppt = Presentation()
        # match your desired pixel size (kept original values)
        self._ppt.slide_width = Pt(1280)
        self._ppt.slide_height = Pt(720)

    def get_sub_element(self, parent, tagname, **kwargs):
        """Helper method to create XML elements"""
        element = OxmlElement(tagname)
        for k, v in kwargs.items():
            element.set(k, v)
        parent.append(element)
        return element

    async def fetch_network_assets(self):
        """
        Download any http(s) images referenced in slides and shapes into temp dir.
        Rewrites shape.picture.path to local file path when downloaded.
        """
        image_urls: List[str] = []
        models_with_network_asset: List[PptxPictureBoxModel] = []

        # collect from global model shapes (if any)
        global_shapes = getattr(self._ppt_model, "shapes", []) or []
        # collect from per-slide shapes
        for shape_model in list(global_shapes):
            if isinstance(shape_model, PptxPictureBoxModel):
                image_path = getattr(shape_model, "picture", None) and getattr(shape_model.picture, "path", "")
                if isinstance(image_path, str) and image_path.startswith("http"):
                    image_urls.append(image_path)
                    models_with_network_asset.append(shape_model)

        for slide in self._slide_models:
            for shape_model in getattr(slide, "shapes", []) or []:
                if isinstance(shape_model, PptxPictureBoxModel):
                    image_path = getattr(shape_model, "picture", None) and getattr(shape_model.picture, "path", "")
                    if isinstance(image_path, str) and image_path.startswith("http"):
                        image_urls.append(image_path)
                        models_with_network_asset.append(shape_model)

        if not image_urls:
            return

        try:
            downloaded_paths = await download_files(image_urls, self._temp_dir)
        except Exception as e:
            print("❌ Failed downloading image URLs:", e)
            downloaded_paths = [None] * len(image_urls)

        # map results back to models (zip will stop at shortest; ensure lengths match)
        for model, path in zip(models_with_network_asset, downloaded_paths):
            if path:
                model.picture.path = path
                model.picture.is_network = False
            else:
                # fallback to placeholder
                model.picture.path = "/static/images/placeholder.jpg"
                model.picture.is_network = False

    async def create_ppt(self):
        """
        Public: fetch remote images then create slides.
        """
        await self.fetch_network_assets()

        for slide_model in self._slide_models:
            # If global shapes exist (list), insert them at start of slide shapes
            global_shapes = getattr(self._ppt_model, "shapes", []) or []
            # ensure we do not mutate the original global list on append
            slide_shapes = list(slide_model.shapes) if getattr(slide_model, "shapes", None) else []
            if isinstance(global_shapes, list):
                # prepend global shapes (copy references)
                slide_shapes = list(global_shapes) + slide_shapes

            # use a shallow copy of slide model with shapes replaced so downstream code uses same interface
            slide_copy = PptxSlideModel(**{**slide_model.__dict__}) if hasattr(PptxSlideModel, "__init__") else slide_model
            # If slide_copy is a dataclass-like Pydantic object, setting attribute is safe; otherwise fallback
            try:
                slide_copy.shapes = slide_shapes
            except Exception:
                slide_model.shapes = slide_shapes

            # call core add/populate
            self.add_and_populate_slide(slide_model)

    def set_presentation_theme(self):
        """
        Optionally replace theme colors if provided in model.
        The model may include a `theme` attribute with color mapping.
        """
        try:
            theme_colors = getattr(self._ppt_model, "theme", {}).get("colors", {}) or {}
            if not theme_colors:
                return
            slide_master = self._ppt.slide_master
            slide_master_part = slide_master.part
            theme_part = slide_master_part.part_related_by(RT.THEME)
            theme = fromstring(theme_part.blob)
            nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
            for color_name, hex_value in theme_colors.items():
                if not color_name:
                    continue
                nodes = theme.xpath(
                    f"a:themeElements/a:clrScheme/a:{color_name}/a:srgbClr",
                    namespaces=nsmap,
                )
                if nodes:
                    color_node = nodes[0]
                    # set hex (without '#') and ensure uppercase
                    hex_clean = hex_value.replace("#", "").upper()
                    color_node.set("val", hex_clean)
            theme_part._blob = tostring(theme)
        except Exception as e:
            print("⚠️ set_presentation_theme failed:", e)

    def add_and_populate_slide(self, slide_model: PptxSlideModel):
        """
        Create a new blank slide and add shapes according to models.
        """
        try:
            slide = self._ppt.slides.add_slide(self._ppt.slide_layouts[BLANK_SLIDE_LAYOUT])
        except Exception as e:
            # fallback: add first layout
            print("⚠️ Could not use BLANK_SLIDE_LAYOUT, falling back:", e)
            slide = self._ppt.slides.add_slide(self._ppt.slide_layouts[0])

        if getattr(slide_model, "background", None):
            try:
                self.apply_fill_to_shape(slide.background, slide_model.background)
            except Exception as e:
                print("⚠️ apply fill failed on background:", e)

        if getattr(slide_model, "note", None):
            try:
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    slide.notes_slide.notes_text_frame.text = slide_model.note
            except Exception:
                # Some PPT layouts may not have notes slide
                pass

        for shape_model in getattr(slide_model, "shapes", []) or []:
            model_type = type(shape_model)
            try:
                if model_type is PptxPictureBoxModel:
                    self.add_picture(slide, shape_model)
                elif model_type is PptxAutoShapeBoxModel:
                    self.add_autoshape(slide, shape_model)
                elif model_type is PptxTextBoxModel:
                    self.add_textbox(slide, shape_model)
                elif model_type is PptxConnectorModel:
                    self.add_connector(slide, shape_model)
                else:
                    # Unknown shape: try text representation
                    try:
                        self.add_textbox(slide, PptxTextBoxModel(position=PptxPositionModel(left=50, top=50, width=400, height=100), paragraphs=[PptxParagraphModel(text=str(shape_model))]))
                    except Exception:
                        print("⚠️ Unknown shape model type:", model_type)
            except Exception as inner:
                print("⚠️ Failed to add shape:", inner)

    def add_connector(self, slide: Slide, connector_model: PptxConnectorModel):
        try:
            if getattr(connector_model, "thickness", 0) == 0:
                return
            connector_shape = slide.shapes.add_connector(
                connector_model.type, *connector_model.position.to_pt_xyxy()
            )
            connector_shape.line.width = Pt(connector_model.thickness)
            connector_shape.line.color.rgb = RGBColor.from_string(connector_model.color)
            self.set_fill_opacity(connector_shape.line.fill, connector_model.opacity)

        except Exception as e:
            print("⚠️ add_connector failed:", e)

    def add_picture(self, slide: Slide, picture_model: PptxPictureBoxModel):
        image_path = getattr(picture_model, "picture", None) and getattr(picture_model.picture, "path", "")
        if not image_path:
            return

        try:
            if (
                getattr(picture_model, "clip", False)
                or getattr(picture_model, "border_radius", None)
                or getattr(picture_model, "invert", False)
                or getattr(picture_model, "opacity", None)
                or getattr(picture_model, "object_fit", None)
                or getattr(picture_model, "shape", None)
            ):
                # open and transform image safely
                image = Image.open(image_path)
                image = image.convert("RGBA")

                if getattr(picture_model, "border_radius", None):
                    image = round_image_corners(image, picture_model.border_radius)

                if getattr(picture_model, "object_fit", None):
                    image = fit_image(
                        image,
                        picture_model.position.width,
                        picture_model.position.height,
                        picture_model.object_fit,
                    )
                elif getattr(picture_model, "clip", False):
                    image = clip_image(
                        image,
                        picture_model.position.width,
                        picture_model.position.height,
                    )

                if getattr(picture_model, "border_radius", None):
                    image = round_image_corners(image, picture_model.border_radius)

                if getattr(picture_model, "shape", None) == PptxBoxShapeEnum.CIRCLE:
                    image = create_circle_image(image)

                if getattr(picture_model, "invert", False):
                    image = invert_image(image)

                if getattr(picture_model, "opacity", None):
                    image = set_image_opacity(image, picture_model.opacity)

                # save transformed image to temp dir
                out_path = os.path.join(self._temp_dir, f"{uuid.uuid4()}.png")
                image.save(out_path)
                image_path = out_path
        except Exception as e:
            print("⚠️ Image processing failed:", e)
            # fallback: if path is remote or invalid, use placeholder
            if not os.path.exists(image_path):
                image_path = "/static/images/placeholder.jpg"

        # compute position with margin
        try:
            margined_position = self.get_margined_position(
                picture_model.position, getattr(picture_model, "margin", None)
            )
            slide.shapes.add_picture(image_path, *margined_position.to_pt_list())
        except Exception as e:
            print("⚠️ Failed to add picture to slide:", e)

    def add_autoshape(self, slide: Slide, autoshape_box_model: PptxAutoShapeBoxModel):
        try:
            position = autoshape_box_model.position
            if getattr(autoshape_box_model, "margin", None):
                position = self.get_margined_position(position, autoshape_box_model.margin)

            autoshape = slide.shapes.add_shape(
                autoshape_box_model.type, *position.to_pt_list()
            )

            textbox = autoshape.text_frame
            textbox.word_wrap = autoshape_box_model.text_wrap

            self.apply_fill_to_shape(autoshape, autoshape_box_model.fill)
            self.apply_margin_to_text_box(textbox, autoshape_box_model.margin)
            self.apply_stroke_to_shape(autoshape, autoshape_box_model.stroke)
            self.apply_shadow_to_shape(autoshape, autoshape_box_model.shadow)
            self.apply_border_radius_to_shape(autoshape, autoshape_box_model.border_radius)

            if getattr(autoshape_box_model, "paragraphs", None):
                self.add_paragraphs(textbox, autoshape_box_model.paragraphs)
        except Exception as e:
            print("⚠️ add_autoshape failed:", e)

    def add_textbox(self, slide: Slide, textbox_model: PptxTextBoxModel):
        try:
            position = textbox_model.position
            textbox_shape = slide.shapes.add_textbox(*position.to_pt_list())
            textbox_shape.width += Pt(2)

            textbox = textbox_shape.text_frame
            textbox.word_wrap = textbox_model.text_wrap

            self.apply_fill_to_shape(textbox_shape, textbox_model.fill)
            self.apply_margin_to_text_box(textbox, textbox_model.margin)
            self.add_paragraphs(textbox, textbox_model.paragraphs)
        except Exception as e:
            print("⚠️ add_textbox failed:", e)

    def add_paragraphs(self, textbox: TextFrame, paragraph_models: List[PptxParagraphModel]):
        for index, paragraph_model in enumerate(paragraph_models):
            paragraph = textbox.add_paragraph() if index > 0 else textbox.paragraphs[0]
            self.populate_paragraph(paragraph, paragraph_model)

    def populate_paragraph(self, paragraph: _Paragraph, paragraph_model: PptxParagraphModel):
        try:
            if paragraph_model.spacing:
                self.apply_spacing_to_paragraph(paragraph, paragraph_model.spacing)

            if paragraph_model.line_height:
                paragraph.line_spacing = paragraph_model.line_height

            if paragraph_model.alignment:
                paragraph.alignment = paragraph_model.alignment

            if paragraph_model.font:
                self.apply_font_to_paragraph(paragraph, paragraph_model.font)

            text_runs = []
            if paragraph_model.text:
                text_runs = self.parse_html_text_to_text_runs(
                    paragraph_model.font, paragraph_model.text
                )
            elif paragraph_model.text_runs:
                text_runs = paragraph_model.text_runs

            for text_run_model in text_runs:
                text_run = paragraph.add_run()
                self.populate_text_run(text_run, text_run_model)
        except Exception as e:
            print("⚠️ populate_paragraph failed:", e)

    def parse_html_text_to_text_runs(self, font: Optional[PptxFontModel], text: str):
        return parse_inline_html_to_runs(text, font)

    def populate_text_run(self, text_run: _Run, text_run_model: PptxTextRunModel):
        try:
            text_run.text = text_run_model.text
            if text_run_model.font:
                self.apply_font(text_run.font, text_run_model.font)
        except Exception as e:
            print("⚠️ populate_text_run failed:", e)

    def apply_border_radius_to_shape(self, shape: Shape, border_radius: Optional[int]):
        if not border_radius:
            return
        try:
            normalized_border_radius = Pt(border_radius) / min(
                max(shape.width, Pt(1)), max(shape.height, Pt(1))
            )
            try:
                shape.adjustments[0] = normalized_border_radius
            except Exception:
                # not all shapes support adjustments
                pass
        except Exception as e:
            print("⚠️ apply_border_radius_to_shape failed:", e)

    def apply_fill_to_shape(self, shape: Shape, fill: Optional[PptxFillModel] = None):
        try:
            if not fill:
                shape.fill.background()
                return
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(fill.color)
            self.set_fill_opacity(shape, fill.opacity)
        except Exception as e:
            print("⚠️ apply_fill_to_shape failed:", e)

    def apply_stroke_to_shape(self, shape: Shape, stroke: Optional[PptxStrokeModel] = None):
        try:
            if not stroke or getattr(stroke, "thickness", 0) == 0:
                shape.line.fill.background()
                return
            shape.line.fill.solid()
            shape.line.fill.fore_color.rgb = RGBColor.from_string(stroke.color)
            shape.line.width = Pt(stroke.thickness)
            self.set_fill_opacity(shape.line.fill, stroke.opacity)
        except Exception as e:
            print("⚠️ apply_stroke_to_shape failed:", e)

    def apply_shadow_to_shape(self, shape: Shape, shadow: Optional[PptxShadowModel] = None):
        """
        Safely apply outerShdw XML. Many PPTX shapes do not support shadow manipulation;
        swallow exceptions and print warnings instead of crashing.
        """
        try:
            sp_element = shape._element
            sp_pr = sp_element.xpath("p:spPr")[0]
            nsmap = sp_pr.nsmap or {"a": "http://schemas.openxmlformats.org/drawingml/2006/main", "p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
            # Ensure effectLst exists
            effect_list = sp_pr.find("a:effectLst", namespaces=nsmap)
            if effect_list is None:
                effect_list = etree.SubElement(sp_pr, f"{{{nsmap['a']}}}effectLst")

            # remove existing outerShdw if present
            old = effect_list.find("a:outerShdw", namespaces=nsmap)
            if old is not None:
                effect_list.remove(old)

            if shadow is None:
                # make a transparent zero shadow
                outer_shadow = etree.SubElement(effect_list, f"{{{nsmap['a']}}}outerShdw", {"blurRad": "0", "dist": "0", "dir": "0"})
                color_element = etree.SubElement(outer_shadow, f"{{{nsmap['a']}}}srgbClr", {"val": "000000"})
                etree.SubElement(color_element, f"{{{nsmap['a']}}}alpha", {"val": "0"})
                return
            # create shadow element with best-effort numeric conversions
            blur = str(int(Pt(shadow.radius))) if getattr(shadow, "radius", None) else "40000"
            offset = str(int(Pt(shadow.offset))) if getattr(shadow, "offset", None) else "40000"
            dir_val = str(int(getattr(shadow, "angle", 0) * 1000))
            outer_shadow = etree.SubElement(
                effect_list,
                f"{{{nsmap['a']}}}outerShdw",
                {
                    "blurRad": blur,
                    "dir": dir_val,
                    "dist": offset,
                    "rotWithShape": "0",
                },
            )
            color_element = etree.SubElement(outer_shadow, f"{{{nsmap['a']}}}srgbClr", {"val": f"{getattr(shadow, 'color', '000000')}"})
            alpha_val = str(int(getattr(shadow, "opacity", 0.0) * 100000)) if getattr(shadow, "opacity", None) is not None else "100000"
            etree.SubElement(color_element, f"{{{nsmap['a']}}}alpha", {"val": alpha_val})
        except Exception as e:
            print("⚠️ apply_shadow_to_shape failed:", e)

    def set_fill_opacity(self, fill, opacity):
        if opacity is None or opacity >= 1.0:
            return
        try:
            alpha = int((opacity) * 100000)
            ts = getattr(fill, "_xPr", None)
            if ts is None:
                return
            sF = ts.solidFill if hasattr(ts, "solidFill") else None
            if sF is None:
                sF = ts.get_or_change_to_srgbClr()
            self.get_sub_element(sF, "a:alpha", val=str(alpha))
        except Exception as e:
            print(f"Could not set fill opacity: {e}")

    def get_margined_position(self, position: PptxPositionModel, margin: Optional[PptxSpacingModel]) -> PptxPositionModel:
        if not margin:
            return position
        left = position.left + margin.left
        top = position.top + margin.top
        width = max(position.width - margin.left - margin.right, 0)
        height = max(position.height - margin.top - margin.bottom, 0)
        return PptxPositionModel(left=left, top=top, width=width, height=height)

    def apply_margin_to_text_box(self, text_frame: TextFrame, margin: Optional[PptxSpacingModel]):
        try:
            text_frame.margin_left = Pt(margin.left if margin else 0)
            text_frame.margin_right = Pt(margin.right if margin else 0)
            text_frame.margin_top = Pt(margin.top if margin else 0)
            text_frame.margin_bottom = Pt(margin.bottom if margin else 0)
        except Exception as e:
            print("⚠️ apply_margin_to_text_box failed:", e)

    def apply_spacing_to_paragraph(self, paragraph: _Paragraph, spacing: PptxSpacingModel):
        try:
            paragraph.space_before = Pt(spacing.top)
            paragraph.space_after = Pt(spacing.bottom)
        except Exception as e:
            print("⚠️ apply_spacing_to_paragraph failed:", e)

    def apply_font_to_paragraph(self, paragraph: _Paragraph, font: PptxFontModel):
        try:
            self.apply_font(paragraph.font, font)
        except Exception as e:
            print("⚠️ apply_font_to_paragraph failed:", e)

    def apply_font(self, font: Font, font_model: PptxFontModel):
        try:
            if font_model.name:
                font.name = font_model.name
            if font_model.color:
                font.color.rgb = RGBColor.from_string(font_model.color)
            font.italic = bool(font_model.italic)
            font.size = Pt(font_model.size) if font_model.size else font.size
            font.bold = bool(getattr(font_model, "font_weight", 400) >= 600)
            if getattr(font_model, "underline", None) is not None:
                font.underline = bool(font_model.underline)
            if getattr(font_model, "strike", None) is not None:
                self.apply_strike_to_font(font, font_model.strike)
        except Exception as e:
            print("⚠️ apply_font failed:", e)

    def apply_strike_to_font(self, font: Font, strike: Optional[bool]):
        try:
            rPr = font._element
            if strike is True:
                rPr.set("strike", "sngStrike")
            elif strike is False:
                rPr.set("strike", "noStrike")
        except Exception as e:
            print("Could not apply strikethrough: ", e)

    def save(self, path: str):
        # Ensure export dir exists
        os.makedirs(os.path.dirname(path), exist_ok=True)
        try:
            self._ppt.save(path)
        except Exception as e:
            print("❌ Failed to save PPTX:", e)
            raise

