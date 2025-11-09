import pdfplumber
from docx import Document
from pptx import Presentation
import os


class DoclingService:
    """
    A safe replacement for the unavailable 'docling' module.
    Converts PDF, DOCX, and PPTX files to plain markdown-like text.
    """

    def __init__(self):
        pass  # No heavy setup needed

    def parse_to_markdown(self, file_path: str) -> str:
        """
        Converts supported documents (PDF, DOCX, PPTX) into markdown text.
        """
        ext = os.path.splitext(file_path)[-1].lower()

        if ext == ".pdf":
            return self._extract_pdf(file_path)
        elif ext == ".docx":
            return self._extract_docx(file_path)
        elif ext == ".pptx":
            return self._extract_pptx(file_path)
        else:
            return f"❌ Unsupported file type: {ext}"

    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF using pdfplumber"""
        text = []
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text.append(page_text)
            return "\n\n".join(text)
        except Exception as e:
            return f"⚠️ PDF extraction failed: {str(e)}"

    def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX using python-docx"""
        try:
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception as e:
            return f"⚠️ DOCX extraction failed: {str(e)}"

    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PPTX using python-pptx"""
        try:
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
            return "\n\n".join(text)
        except Exception as e:
            return f"⚠️ PPTX extraction failed: {str(e)}"
